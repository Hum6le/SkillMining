from __future__ import annotations

import re
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / ".pptx_deps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent / "skill_induction_proposal_zh_final.pptx"
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

NAVY=RGBColor(20,37,64); INK=RGBColor(42,53,69); MUTED=RGBColor(101,116,133)
BLUE=RGBColor(39,104,180); TEAL=RGBColor(17,139,129); ORANGE=RGBColor(220,121,43)
RED=RGBColor(187,63,69); GREEN=RGBColor(49,138,90); WHITE=RGBColor(255,255,255)
GRID=RGBColor(218,226,236); LIGHT=RGBColor(246,248,251); PB=RGBColor(232,242,253)
PT=RGBColor(229,246,242); PO=RGBColor(253,240,224); PR=RGBColor(252,234,235)
FONT="Microsoft YaHei"

SKILL_MINING_ROOT = Path(r"D:\paper\Skill Mining")
AWM_SKILL_PATH = SKILL_MINING_ROOT / "skill-awm.md"
HG_SKILL_PATH = SKILL_MINING_ROOT / "skill-hg.md"
TRACE_SKILL_PATH = SKILL_MINING_ROOT / "skill-Trace2Skill.md"


def skill_section(path: Path, start: str, end: str, max_chars: int = 720) -> str:
    """Read a bounded excerpt from a generated skill for evidence slides."""
    source = path.read_text(encoding="utf-8")
    pattern = re.escape(start) + r"(.*?)(?=" + re.escape(end) + r"|\Z)"
    match = re.search(pattern, source, flags=re.DOTALL)
    excerpt = (match.group(0) if match else source).strip()
    if len(excerpt) > max_chars:
        excerpt = excerpt[:max_chars].rstrip() + "\n..."
    return excerpt


INSTANCE_ACTION_EXCERPT = skill_section(
    HG_SKILL_PATH, "2. **Verification Method Selection", "3. **Post-Verification Action"
)
STATE_WORKFLOW_EXCERPT = skill_section(
    TRACE_SKILL_PATH, "## Action selection", "## Slot handling"
)


def box(slide,x,y,w,h,fill=WHITE,line=None,rounded=False):
    s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill; s.line.color.rgb=line or fill
    if rounded: s.adjustments[0]=0.08
    return s

def text(slide,s,x,y,w,h,size=18,color=INK,bold=False,align=PP_ALIGN.LEFT,valign=MSO_ANCHOR.TOP,font=FONT):
    b=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=Inches(.06); tf.margin_right=Inches(.06); tf.margin_top=Inches(.04); tf.margin_bottom=Inches(.04); tf.vertical_anchor=valign
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=s; r.font.name=font; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return b

def bullets(slide,items,x,y,w,h,size=17,color=INK):
    b=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.clear(); tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text="• "+item; p.space_after=Pt(8); p.font.name=FONT; p.font.size=Pt(size); p.font.color.rgb=color
    return b

def line(slide,x1,y1,x2,y2,color=GRID,width=1.0):
    c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2)); c.line.color.rgb=color; c.line.width=Pt(width); return c

def header(slide,section,title,n):
    text(slide,section.upper(),.58,.28,4.8,.25,10,BLUE,True)
    text(slide,title,.55,.62,11.8,.58,27,NAVY,True)
    text(slide,f"{n:02d}",12.25,.32,.5,.25,10,MUTED,True,PP_ALIGN.RIGHT)
    line(slide,.58,1.35,12.75,1.35,GRID,.8)

def footer(slide,n):
    line(slide,.58,7.12,12.75,7.12,GRID,.7)
    text(slide,"从轨迹中学习可复用、可分支、可验证的 Skill",.6,7.19,6.0,.17,8,MUTED)
    text(slide,str(n),12.2,7.19,.42,.17,8,MUTED,False,PP_ALIGN.RIGHT)

def card(slide,x,y,w,h,title,body,fill=WHITE,accent=BLUE,size=16):
    box(slide,x,y,w,h,fill,GRID,True); box(slide,x,y,.08,h,accent,accent,True)
    text(slide,title,x+.25,y+.18,w-.45,.32,16,NAVY,True)
    text(slide,body,x+.25,y+.62,w-.45,h-.78,size,INK)

def pill(slide,s,x,y,w,fill,color=INK):
    box(slide,x,y,w,.34,fill,fill,True); text(slide,s,x+.04,y+.04,w-.08,.23,10,color,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)

def arrow(slide,x1,y,x2,color=BLUE):
    line(slide,x1,y,x2,y,color,2); text(slide,">",(x1+x2)/2-.1,y-.17,.2,.3,17,color,True,PP_ALIGN.CENTER)

def workflow(slide,labels,x,y,widths,fills):
    cur=x
    for i,label in enumerate(labels):
        box(slide,cur,y,widths[i],.8,fills[i],fills[i],True); text(slide,label,cur+.06,y+.14,widths[i]-.12,.5,14,NAVY,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
        if i<len(labels)-1: arrow(slide,cur+widths[i]+.03,y+.4,cur+widths[i]+.38)
        cur+=widths[i]+.43

# 1
s=prs.slides.add_slide(BLANK); box(s,0,0,13.333,7.5,NAVY,NAVY); box(s,0,5.98,13.333,1.52,BLUE,BLUE)
text(s,"研究方案",.72,.72,3.0,.28,12,RGBColor(168,211,248),True)
text(s,"从轨迹中学习\n可复用、可分支的 Skill",.7,1.45,10.5,1.25,36,WHITE,True)
text(s,"两个方向：Macro Skill  ·  State-conditioned Skill",.75,3.38,8.8,.36,20,RGBColor(218,232,248))
text(s,"作者姓名 · 单位",.75,6.38,5,.28,15,WHITE,True)
text(s,"Research Proposal",10.0,6.4,2.4,.25,12,RGBColor(218,232,248),False,PP_ALIGN.RIGHT)

# 2 overall
s=prs.slides.add_slide(BLANK); header(s,"研究动机","现有方法能生成 Skill，但还没有解决两个结构性问题",2)
text(s,"AWM 与 Trace2Skill 的输出分别暴露出两种问题：一个偏向记忆具体轨迹，一个缺少状态条件。",.85,1.68,11.6,.48,20,NAVY,True,PP_ALIGN.CENTER)
card(s,.75,2.45,5.55,2.75,"方向一：泛化性不足","Skill 记住了具体姓名、账号、slot 值或动作序列。\n\n问题：它描述的是一条轨迹，而不是可替换参数的程序。",PO,ORANGE,18)
card(s,7.0,2.45,5.55,2.75,"方向二：分支条件缺失","不同状态下都出现了局部规则。\n\n问题：规则被平铺后，看起来互相冲突，实际缺少状态条件。",PB,BLUE,18)
text(s,"共同缺口：轨迹经验没有被组织成可复用、带条件的程序性知识。",1.0,6.15,11.3,.4,21,TEAL,True,PP_ALIGN.CENTER); footer(s,2)

# Replace the generic overview with direct AWM source evidence.
box(s,.42,1.5,12.5,5.45,WHITE,WHITE,False)
text(s,"AWM 原文证据：同一任务中的验证条件并未被统一",.78,1.72,8.2,.32,20,NAVY,True)
card(s,.72,2.25,5.72,3.55,"原文片段 A · Standard Recovery Flow","“Ensure at least two distinct identifiers are provided”\n\n→ 标准流程要求两个身份标识。\n\n来源：skill-awm.md",PO,ORANGE,17)
card(s,6.9,2.25,5.72,3.55,"原文片段 B · Identity Verification Priorities","“Secondary: Request one of the following”\n\n→ 决策规则又描述为请求一个辅助标识。\n\n来源：skill-awm.md",PB,BLUE,17)
pill(s,"问题不是简单选一条规则，而是缺少区分状态的 guard",3.25,6.12,6.85,PR,RED)

# 3 direction 1 problem
s=prs.slides.add_slide(BLANK); header(s,"方向一 · Macro Skill","真实 Skill 原文已经暴露了实例记忆问题",3)
text(s,"原文来源：Skill Mining/skill-hg.md",.85,1.68,7.2,.28,12,ORANGE,True)
box(s,.75,2.05,7.0,3.75,PO,PO,True)
text(s,INSTANCE_ACTION_EXCERPT,.98,2.3,6.55,3.2,15,INK,False,font="Consolas")
card(s,8.15,2.05,4.2,3.75,"问题在哪里？","动作序列直接绑定了：\n\n• aphoenix1\n• cm374950\n• alessandro phoenix\n• 57820\n• (499) 412-7409\n\n这些是训练实例，不是可替换变量。",PR,RED,16)
pill(s,"原文证据：instance-specific actions",1.65,6.18,4.5,PR,RED)
text(s,"应当学习 enter-details(PERSON / ACCOUNT / CREDENTIAL)，而不是记住某几条具体轨迹。",6.55,6.2,5.9,.38,16,TEAL,True,PP_ALIGN.CENTER)
footer(s,3)

# 4 direction1 solution
s=prs.slides.add_slide(BLANK); header(s,"方向一 · 解决方案","从 Trace Skill 到 Macro Skill：学习可组合的程序模板",4)
card(s,.75,1.8,3.65,3.8,"1. 抽取结构","从轨迹中识别：\n\n状态、动作、slot、结果\n\n保留原始值用于追溯，\n使用类型用于泛化。",PB,BLUE,17)
card(s,4.85,1.8,3.65,3.8,"2. 变量化规则","将具体实例归一化为：\n\nPERSON_NAME\nPHONE\nACCOUNT_ID\n\n形成参数化动作模板。",PT,TEAL,17)
card(s,8.95,1.8,3.65,3.8,"3. 跨轨迹验证","在新实体、新措辞、\n新 slot 组合上回放。\n\n只有可迁移的结构\n才能提升为 Macro Skill。",PO,ORANGE,17)
text(s,"具体轨迹  →  参数化子程序  →  可迁移 Macro Skill",1.1,6.25,11.0,.35,21,NAVY,True,PP_ALIGN.CENTER); footer(s,4)

# 5 direction1 solution
s=prs.slides.add_slide(BLANK); header(s,"方向一 · 解决方案","从 Trace Skill 到 Macro Skill：学习可组合的程序模板",5)
card(s,.75,1.8,3.65,3.8,"1. 抽取结构","从轨迹中识别：\n\n状态、动作、slot、结果\n\n保留原始值用于追溯，\n使用类型用于泛化。",PB,BLUE,17)
card(s,4.85,1.8,3.65,3.8,"2. 变量化规则","将具体实例归一化为：\n\nPERSON_NAME\nPHONE\nACCOUNT_ID\n\n形成参数化动作模板。",PT,TEAL,17)
card(s,8.95,1.8,3.65,3.8,"3. 跨轨迹验证","在新实体、新措辞、\n新 slot 组合上回放。\n\n只有可迁移的结构\n才能提升为 Macro Skill。",PO,ORANGE,17)
text(s,"具体轨迹  →  参数化子程序  →  可迁移 Macro Skill",1.1,6.25,11.0,.35,21,NAVY,True,PP_ALIGN.CENTER); footer(s,5)

# 6 direction 2 problem
s=prs.slides.add_slide(BLANK); header(s,"方向二 · State-conditioned Skill","真实 Skill 有状态描述，但关键 guard 仍不够可执行",6)
text(s,"原文来源：Skill Mining/skill-Trace2Skill.md",.85,1.62,10.7,.28,12,BLUE,True)
box(s,.75,2.0,7.15,3.95,PB,PB,True)
text(s,STATE_WORKFLOW_EXCERPT,.98,2.23,6.68,3.5,14,INK,False,font="Consolas")
box(s,8.15,2.0,4.2,3.95,PR,PR,True)
text(s,"问题在哪里？",8.48,2.3,2.6,.3,15,RED,True)
bullets(s,["“请求两个凭证”没有明确当前已收集哪些字段。","没有把 credential count、字段可用性和系统返回结果写成可判定 guard。","同一段文字同时描述 ask、verify、失败重试，执行边界仍依赖模型临场推断。"],8.42,2.78,3.55,2.15,14,INK)
text(s,"需要 State–Guard–Action–Transition，而不仅是流程叙述。",8.42,5.22,3.55,.45,16,RED,True,PP_ALIGN.CENTER)
footer(s,6)

# 7 direction2 solution
s=prs.slides.add_slide(BLANK); header(s,"方向二 · 解决方案","把规则表示成 State–Guard–Action–Transition",7)
workflow(s,["State\n当前状态","Guard\n前置条件","Action\n执行动作","Transition\n下一状态"],.85,2.15,[2.0,2.0,1.85,2.15],[PB,PO,PT,PB])
text(s,"示例",.9,3.65,1.0,.28,12,TEAL,True)
box(s,.85,4.02,11.65,1.15,PT,PT,True)
text(s,"account_identified + {name, phone, email}  →  verify-identity  →  identity_verified",1.15,4.36,11.0,.35,21,NAVY,True,PP_ALIGN.CENTER,font="Consolas")
text(s,"account_identified + {name, phone}  →  ask-for-credential  →  verification_incomplete",1.15,5.02,11.0,.3,16,INK,False,PP_ALIGN.CENTER,font="Consolas")
pill(s,"规则不是互相覆盖，而是由状态条件选择",3.8,5.9,5.75,PT,TEAL)
footer(s,7)

# 8 unified mechanism
s=prs.slides.add_slide(BLANK); header(s,"统一机制","Reference-use policy：教会 Skill 何时查、查什么、如何用",8)
workflow(s,["观察状态","判断是否需要查","检索对应证据","ground + apply","验证后执行"],.65,2.15,[1.5,1.75,1.7,1.65,1.55],[PB,PO,PT,PB,PT])
card(s,.85,3.85,3.55,1.7,"何时查？","状态不确定、动作边界不清、slot schema 不确定时。",PO,ORANGE,16)
card(s,4.9,3.85,3.55,1.7,"查什么？","查 action schema、branch condition、boundary case，而非只查关键词。",PT,TEAL,16)
card(s,8.95,3.85,3.55,1.7,"如何用？","把 reference 当规则证据；slot 值必须来自当前对话。",PB,BLUE,16)
text(s,"Reference 不再是被动示例，而是 Skill 执行过程中的知识工具。",1.0,6.25,11.3,.35,20,NAVY,True,PP_ALIGN.CENTER); footer(s,8)

# 9 model
s=prs.slides.add_slide(BLANK); header(s,"方法框架","两类 Skill 共同构成一个可执行的表示",9)
card(s,.78,1.85,5.55,3.55,"Task Policy","state → business action\n\naccount_identified + enough credentials\n→ verify-identity\n\nidentity_verified\n→ send-link\n\nverification_failed\n→ retry / escalate",PB,BLUE,17)
card(s,7.0,1.85,5.55,3.55,"Reference-use Policy","state + uncertainty → retrieval\n\nretrieve-or-act\nquery construction\nevidence grounding\napplicability verification",PT,TEAL,17)
text(s,"Skill = π_task ⊕ π_reference",1.0,5.95,11.3,.48,28,NAVY,True,PP_ALIGN.CENTER,font="Consolas")
text(s,"前者决定做什么，后者决定何时查证据以及如何把证据转成行动。",1.0,6.48,11.3,.28,15,MUTED,False,PP_ALIGN.CENTER); footer(s,9)

# 10 evaluation
s=prs.slides.add_slide(BLANK); header(s,"实验计划","用 bad case 验证两条方向，而不是只看总体分数",10)
items=[("Macro 泛化","训练中出现流程，测试中替换实体 / 说法 / slot 组合","entity grounding · unseen-variant success",PO,ORANGE),
       ("分支一致性","相同 intent，控制当前 state 和已有凭证数量","branch accuracy · conflict rate",PT,TEAL),
       ("Reference 使用","对比不查、盲查、state-aware 查阅","retrieval value · misuse rate · token efficiency",PB,BLUE)]
for i,(a,b,c,fill,accent) in enumerate(items):
    y=1.78+i*1.25; box(s,.75,y,2.15,.88,fill,fill,True); text(s,a,.92,y+.25,1.8,.3,16,accent,True,PP_ALIGN.CENTER)
    box(s,3.15,y,4.7,.88,LIGHT,GRID,True); text(s,b,3.38,y+.18,4.25,.5,15,NAVY,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    box(s,8.1,y,4.35,.88,WHITE,GRID,True); text(s,c,8.35,y+.18,3.85,.5,15,INK,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
pill(s,"问题页已使用真实生成 Skill 原文摘录",3.45,5.75,6.4,PT,TEAL)
text(s,"最小证据单元：原始轨迹 + 生成 Skill + reference + 预测 + 结果。",1.0,6.35,11.3,.3,16,MUTED,False,PP_ALIGN.CENTER); footer(s,10)

# 11 contributions
s=prs.slides.add_slide(BLANK); header(s,"贡献","我们的贡献不是更长的 Skill，而是更会使用证据的 Skill",11)
box(s,.8,1.75,11.75,1.4,NAVY,NAVY,True); text(s,"从具体轨迹中学习可迁移的 Macro Skill，\n并用状态条件组织不同分支规则。",1.15,2.08,11.05,.75,25,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
bullets(s,["Macro Skill：把实例行为提升为可替换、可迁移的程序结构","State-conditioned Skill：用 State–Guard–Action–Transition 区分分支","Reference-use Policy：学习何时查、查什么、如何 grounding 和验证","Bad-case-driven evaluation：用可控反例验证泛化与规则一致性"],1.0,3.75,10.0,1.8,18)
box(s,9.95,4.05,2.2,1.35,PT,PT,True); text(s,"可复用\n可分支\n可验证",10.18,4.28,1.75,.85,21,TEAL,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
text(s,"Thank you",10.0,6.3,2.3,.3,20,BLUE,True,PP_ALIGN.CENTER); footer(s,11)

prs.save(OUT); print(f"saved {OUT}")
