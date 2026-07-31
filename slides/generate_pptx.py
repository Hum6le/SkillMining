from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / ".pptx_deps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent / "skill_reference_story.pptx"

NAVY = RGBColor(24, 36, 58)
INK = RGBColor(37, 47, 63)
MUTED = RGBColor(91, 105, 124)
BLUE = RGBColor(39, 104, 180)
TEAL = RGBColor(19, 142, 133)
ORANGE = RGBColor(224, 126, 54)
RED = RGBColor(191, 67, 67)
GREEN = RGBColor(45, 137, 89)
PALE_BLUE = RGBColor(232, 242, 252)
PALE_TEAL = RGBColor(228, 246, 242)
PALE_ORANGE = RGBColor(252, 239, 225)
PALE_RED = RGBColor(251, 232, 232)
PALE_GRAY = RGBColor(244, 247, 250)
WHITE = RGBColor(255, 255, 255)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill=WHITE, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def line(slide, x1, y1, x2, y2, color=MUTED, width=1.5):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def text(slide, value, x, y, w, h, size=18, color=INK, bold=False,
         align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP,
         margin=0.06):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=18, color=INK, gap=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.04)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "- " + item
        p.level = 0
        p.space_after = Pt(gap)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def title(slide, kicker, heading, number=None):
    text(slide, kicker.upper(), 0.58, 0.35, 4.8, 0.28, 10, BLUE, True)
    text(slide, heading, 0.55, 0.72, 11.7, 0.58, 27, NAVY, True)
    if number is not None:
        text(slide, f"{number:02d}", 12.35, 0.38, 0.45, 0.3, 11, MUTED, True, PP_ALIGN.RIGHT)
    line(slide, 0.58, 1.42, 12.75, 1.42, RGBColor(218, 225, 234), 0.8)


def footer(slide, n):
    line(slide, 0.58, 7.16, 12.75, 7.16, RGBColor(225, 230, 236), 0.7)
    text(slide, "Trajectory-to-Skill Research Direction", 0.6, 7.22, 5.5, 0.18, 8, MUTED)
    text(slide, str(n), 12.2, 7.22, 0.45, 0.18, 8, MUTED, False, PP_ALIGN.RIGHT)


def add_note(slide, note):
    # python-pptx does not expose notes authoring consistently across versions;
    # notes are also saved in slides/speaker_notes.md.
    return None


def pill(slide, label, x, y, w, fill, color=INK):
    rect(slide, x, y, w, 0.34, fill, fill, True)
    text(slide, label, x + 0.04, y + 0.02, w - 0.08, 0.25, 10, color, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def card(slide, x, y, w, h, heading, body, fill=WHITE, accent=BLUE, body_size=15):
    rect(slide, x, y, w, h, fill, RGBColor(221, 228, 236), True)
    rect(slide, x, y, 0.08, h, accent, accent, True)
    text(slide, heading, x + 0.25, y + 0.2, w - 0.45, 0.34, 16, NAVY, True)
    text(slide, body, x + 0.25, y + 0.67, w - 0.45, h - 0.83, body_size, INK)


def workflow(slide, nodes, y=3.0, x=0.75, widths=None, colors=None):
    widths = widths or [1.7] * len(nodes)
    colors = colors or [PALE_BLUE] * len(nodes)
    cur = x
    for i, node in enumerate(nodes):
        rect(slide, cur, y, widths[i], 0.82, colors[i], colors[i], True)
        text(slide, node, cur + 0.07, y + 0.16, widths[i] - 0.14, 0.5, 14, NAVY, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(nodes) - 1:
            line(slide, cur + widths[i], y + 0.41, cur + widths[i] + 0.43, y + 0.41, BLUE, 2.0)
            text(slide, ">", cur + widths[i] + 0.12, y + 0.22, 0.18, 0.28, 18, BLUE, True, PP_ALIGN.CENTER)
        cur += widths[i] + 0.43


# 1: title
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, 13.333, 7.5, NAVY, NAVY)
rect(slide, 0, 5.95, 13.333, 1.55, BLUE, BLUE)
text(slide, "RESEARCH DIRECTION", 0.72, 0.72, 4.0, 0.3, 12, RGBColor(163, 207, 247), True)
text(slide, "From Trajectory Traces\nto Reusable Skills", 0.7, 1.45, 8.9, 1.45, 34, WHITE, True)
text(slide, "Bad case  →  cause  →  insight", 0.75, 3.35, 6.4, 0.38, 20, RGBColor(214, 230, 247), False)
text(slide, "Two directions: macro skills and reference-use skills", 0.75, 6.35, 8.5, 0.4, 18, WHITE, True)
text(slide, "Working deck · evidence gaps explicitly marked", 9.1, 6.37, 3.45, 0.34, 11, RGBColor(214, 230, 247), False, PP_ALIGN.RIGHT)

# 2
slide = prs.slides.add_slide(blank); title(slide, "Starting point", "We have generated skills — but not yet a causal story", 2)
card(slide, 0.7, 1.85, 3.75, 3.8, "Observed locally", "Generated artifacts exist:\n\n• AWM-style workflows\n• Trace2Skill action protocol\n• HG error analysis\n• reference.md files", PALE_TEAL, TEAL)
card(slide, 4.78, 1.85, 3.75, 3.8, "Missing currently", "No aligned AWM-vs-Trace2Skill\nper-turn failure table.\n\nTherefore: candidate bad cases\nare marked, not claimed as results.", PALE_ORANGE, ORANGE)
card(slide, 8.86, 1.85, 3.75, 3.8, "Research need", "Turn artifacts into a testable story:\n\nBad case\n→ root cause\n→ design insight\n→ measurable intervention", PALE_BLUE, BLUE)
text(slide, "Evidence discipline: do not turn a plausible failure into an empirical claim.", 1.05, 6.1, 11.25, 0.45, 20, NAVY, True, PP_ALIGN.CENTER)
footer(slide, 2)

# 3
slide = prs.slides.add_slide(blank); title(slide, "Generated artifacts", "Two skills expose complementary symptoms", 3)
card(slide, 0.72, 1.78, 5.7, 4.55, "AWM-style workflow", "Rich local advice\n\n• many branches and tone rules\n• concrete names, IDs, and slot values\n• “follow the action sequence”\n• rules can disagree on conditions", PALE_ORANGE, ORANGE, 17)
card(slide, 6.92, 1.78, 5.7, 4.55, "Trace2Skill protocol", "Precise local constraints\n\n• action-first\n• ordered slots\n• no stale values\n• but state / action schema stay implicit", PALE_BLUE, BLUE, 17)
pill(slide, "rich but instance-heavy", 1.35, 5.68, 2.65, PALE_ORANGE, ORANGE)
pill(slide, "precise but under-contextualized", 8.15, 5.68, 3.25, PALE_BLUE, BLUE)
text(slide, "Common question: how does the skill know what evidence to consult, and how to apply it now?", 1.0, 6.55, 11.3, 0.35, 16, NAVY, True, PP_ALIGN.CENTER)
footer(slide, 3)

# 4
slide = prs.slides.add_slide(blank); title(slide, "Candidate bad case A", "Macro skill fails when a workflow memorizes instances", 4)
text(slide, "Observed artifact", 0.8, 1.78, 2.0, 0.3, 12, ORANGE, True)
rect(slide, 0.75, 2.15, 5.45, 2.35, PALE_ORANGE, PALE_ORANGE, True)
text(slide, "verify-identity:chloe zhang,83942,(482) ...\nverify-identity:albert sanders,86227,...\npull-up-account:crystal minh\n\nFollow the action sequence above.", 1.03, 2.48, 4.9, 1.7, 18, INK, False, font="Consolas")
text(slide, "Candidate test", 6.72, 1.78, 2.0, 0.3, 12, BLUE, True)
rect(slide, 6.65, 2.15, 5.85, 2.35, PALE_BLUE, PALE_BLUE, True)
workflow(slide, ["known\nworkflow", "new\nentity", "same\nprocedure"], y=2.57, x=7.03, widths=[1.55, 1.55, 1.55], colors=[PALE_BLUE, WHITE, PALE_TEAL])
text(slide, "Test whether concrete values are replaced by variables,\nnot copied from training traces.", 7.02, 3.72, 5.1, 0.55, 16, NAVY, True, PP_ALIGN.CENTER)
pill(slide, "TO FILL: paired unseen-entity case", 3.65, 5.28, 5.9, PALE_RED, RED)
text(slide, "Root cause hypothesis: trajectory examples are stored as procedures, but not parameterized as reusable operators.", 1.0, 6.18, 11.3, 0.48, 17, NAVY, True, PP_ALIGN.CENTER)
footer(slide, 4)

# 5
slide = prs.slides.add_slide(blank); title(slide, "Candidate bad case B", "Rules look contradictory when state conditions are flattened", 5)
text(slide, "Two locally reasonable rules", 0.8, 1.75, 3.0, 0.3, 12, RED, True)
card(slide, 0.75, 2.15, 5.45, 1.35, "Rule A", "Two valid identifiers → verify-identity", PALE_BLUE, BLUE, 17)
card(slide, 0.75, 3.78, 5.45, 1.35, "Rule B", "Request one secondary identifier", PALE_ORANGE, ORANGE, 17)
text(slide, "Missing variable", 6.75, 1.75, 2.4, 0.3, 12, TEAL, True)
rect(slide, 6.65, 2.15, 5.75, 2.98, PALE_TEAL, PALE_TEAL, True)
text(slide, "current state", 7.0, 2.45, 1.6, 0.3, 13, MUTED, True)
text(slide, "account identified\ncredentials = {name, phone}\nidentity not verified", 7.0, 2.82, 4.7, 0.78, 19, NAVY, True)
line(slide, 7.55, 3.85, 7.55, 4.22, TEAL, 2.2)
text(slide, "state-conditioned branch", 7.0, 4.32, 2.3, 0.28, 13, TEAL, True)
text(slide, "if enough credentials → verify\nelse → ask for missing credential", 7.0, 4.67, 4.9, 0.58, 18, NAVY, True)
pill(slide, "not conflict resolution alone", 3.7, 5.65, 2.9, PALE_RED, RED)
pill(slide, "branch-condition learning", 6.8, 5.65, 2.9, PALE_TEAL, TEAL)
text(slide, "TO FILL: same-state / different-state paired trajectory case", 2.35, 6.35, 8.7, 0.3, 13, MUTED, False, PP_ALIGN.CENTER)
footer(slide, 5)

# 6
slide = prs.slides.add_slide(blank); title(slide, "Diagnosis", "The missing object is a reference-use policy", 6)
workflow(slide, ["dialogue\nstate", "retrieve?", "reference\nevidence", "apply\nrule", "act"], y=2.45, x=0.85, widths=[1.7, 1.35, 1.85, 1.45, 1.25], colors=[PALE_BLUE, PALE_ORANGE, PALE_TEAL, PALE_BLUE, PALE_TEAL])
text(slide, "Current implicit behavior", 1.05, 1.78, 3.0, 0.3, 12, RED, True)
text(slide, "query by keywords → paste snippets → generate", 1.05, 3.72, 5.2, 0.42, 19, RED, True)
text(slide, "Proposed object", 7.15, 1.78, 2.0, 0.3, 12, TEAL, True)
text(slide, "a skill that knows when\nwhat and how to retrieve", 7.15, 3.45, 4.75, 0.85, 26, NAVY, True)
text(slide, "Reference is not passive context.\nIt is part of the procedure.", 7.18, 4.58, 4.7, 0.6, 18, TEAL, True)
footer(slide, 6)

# 7
slide = prs.slides.add_slide(blank); title(slide, "Core representation", "Skill = task policy + reference-use policy", 7)
card(slide, 0.8, 1.9, 5.5, 3.8, "Task policy", "state → business action\n\nExamples:\n• account_identified → verify\n• identity_verified → send-link\n• verification_failed → retry / escalate", PALE_BLUE, BLUE, 17)
card(slide, 7.0, 1.9, 5.5, 3.8, "Reference-use policy", "state + uncertainty → query → evidence use\n\nExamples:\n• retrieve schema\n• retrieve boundary case\n• ground current slot values\n• validate before acting", PALE_TEAL, TEAL, 17)
line(slide, 6.3, 3.75, 7.0, 3.75, NAVY, 2.5)
text(slide, "+", 6.44, 3.53, 0.34, 0.36, 24, NAVY, True, PP_ALIGN.CENTER)
text(slide, "A reusable skill must encode both decisions.", 1.2, 6.25, 10.9, 0.4, 20, NAVY, True, PP_ALIGN.CENTER)
footer(slide, 7)

# 8
slide = prs.slides.add_slide(blank); title(slide, "Proposed loop", "Retrieve, ground, apply, verify", 8)
workflow(slide, ["observe\nstate", "decide\nretrieve", "query\nreference", "ground\ncurrent values", "apply +\nverify", "act"], y=2.45, x=0.62, widths=[1.4, 1.45, 1.45, 1.65, 1.5, 1.0], colors=[PALE_BLUE, PALE_ORANGE, PALE_TEAL, PALE_BLUE, PALE_ORANGE, PALE_TEAL])
text(slide, "Reference-use checklist", 0.9, 4.25, 2.5, 0.3, 12, TEAL, True)
bullets(slide, ["Is the current state uncertain?", "What evidence is needed: schema, branch, or example?", "Which values come from the current dialogue?", "Does the evidence apply to this state?"], 0.9, 4.65, 5.7, 1.45, 16)
rect(slide, 7.15, 4.25, 5.2, 1.72, PALE_ORANGE, PALE_ORANGE, True)
text(slide, "Key distinction", 7.5, 4.55, 2.0, 0.28, 13, ORANGE, True)
text(slide, "Examples are evidence —\nnot values to copy.", 7.5, 4.95, 4.3, 0.6, 24, NAVY, True)
footer(slide, 8)

# 9
slide = prs.slides.add_slide(blank); title(slide, "Learning target", "Teach the skill three reference decisions", 9)
card(slide, 0.72, 1.9, 3.75, 3.8, "1 · Retrieve or act", "Given state s and uncertainty u:\n\nShould the agent act now, or consult reference first?", PALE_BLUE, BLUE, 17)
card(slide, 4.8, 1.9, 3.75, 3.8, "2 · Query what matters", "Query by state and decision need:\n\nAction schema?\nBoundary condition?\nRecovery branch?", PALE_TEAL, TEAL, 17)
card(slide, 8.88, 1.9, 3.75, 3.8, "3 · Use and verify", "Separate:\n\n• hard constraints\n• procedural examples\n• current slot values\n• irrelevant instance details", PALE_ORANGE, ORANGE, 17)
text(slide, "Training trace: state → query → evidence → action → outcome", 1.25, 6.25, 10.8, 0.4, 20, NAVY, True, PP_ALIGN.CENTER)
footer(slide, 9)

# 10
slide = prs.slides.add_slide(blank); title(slide, "Evidence plan", "What bad cases do we still need?", 10)
headers = [("Case", "Construct", "Measure"), ("Macro", "Seen workflow + unseen entities", "parameterized generalization"), ("Branch", "Same task + different state", "conditional action accuracy"), ("Reference", "Retrieved example conflicts with dialogue", "grounding / misuse rate")]
x0, y0 = 0.72, 1.9
widths = [1.65, 5.45, 4.9]
cur = x0
for j, h in enumerate(headers[0]):
    rect(slide, cur, y0, widths[j], 0.58, NAVY, NAVY, False)
    text(slide, h, cur + 0.12, y0 + 0.13, widths[j] - 0.24, 0.3, 14, WHITE, True)
    cur += widths[j] + 0.05
for i, row in enumerate(headers[1:], start=1):
    cur = x0
    fill = PALE_GRAY if i % 2 else WHITE
    for j, value in enumerate(row):
        rect(slide, cur, y0 + i * 0.72, widths[j], 0.62, fill, RGBColor(221, 228, 236), False)
        text(slide, value, cur + 0.12, y0 + i * 0.72 + 0.13, widths[j] - 0.24, 0.35, 14, INK, j == 0)
        cur += widths[j] + 0.05
pill(slide, "Current status: case table not yet populated", 3.65, 5.0, 5.95, PALE_RED, RED)
text(slide, "Required artifact: same test trajectory, generated skill, retrieved reference, prediction, outcome.", 1.0, 5.75, 11.3, 0.4, 17, NAVY, True, PP_ALIGN.CENTER)
text(slide, "No fabricated numbers in this deck.", 1.0, 6.35, 11.3, 0.28, 13, MUTED, False, PP_ALIGN.CENTER)
footer(slide, 10)

# 11
slide = prs.slides.add_slide(blank); title(slide, "Takeaway", "The research question is not “how to write longer skills”", 11)
rect(slide, 0.82, 1.95, 11.72, 1.55, NAVY, NAVY, True)
text(slide, "When uncertain, can a skill retrieve the right evidence\nand turn it into the right action?", 1.2, 2.28, 10.95, 0.8, 28, WHITE, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
text(slide, "Contribution draft", 0.9, 4.1, 2.2, 0.3, 12, BLUE, True)
bullets(slide, ["Macro skills: parameterize reusable workflow structure", "Branch-aware rules: condition on state, not flat text", "Reference-use policy: retrieve, ground, apply, verify", "Evaluation: bad case → cause → insight → measurable fix"], 1.0, 4.5, 8.8, 1.45, 19)
pill(slide, "Next: collect paired bad cases", 9.2, 5.1, 2.8, PALE_ORANGE, ORANGE)
footer(slide, 11)


prs.save(OUT)
print(f"saved {OUT}")
