from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / ".pptx_deps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent / "external_research_deck.pptx"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

NAVY = RGBColor(19, 36, 64)
INK = RGBColor(39, 50, 66)
MUTED = RGBColor(100, 115, 133)
BLUE = RGBColor(33, 104, 181)
TEAL = RGBColor(15, 139, 130)
ORANGE = RGBColor(222, 119, 45)
RED = RGBColor(186, 63, 69)
GREEN = RGBColor(47, 139, 91)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(246, 248, 251)
PALE_BLUE = RGBColor(232, 242, 253)
PALE_TEAL = RGBColor(229, 246, 242)
PALE_ORANGE = RGBColor(253, 240, 224)
PALE_RED = RGBColor(252, 234, 235)
GRID = RGBColor(218, 226, 236)


def shape(slide, x, y, w, h, fill=WHITE, line=None, rounded=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill
    if rounded:
        s.adjustments[0] = 0.08
    return s


def tx(slide, s, x, y, w, h, size=18, color=INK, bold=False,
       align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(.06); tf.margin_right = Inches(.06)
    tf.margin_top = Inches(.04); tf.margin_bottom = Inches(.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s; r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    return box


def bullet_box(slide, items, x, y, w, h, size=17, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item; p.space_after = Pt(8)
        p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = color
    return box


def ln(slide, x1, y1, x2, y2, color=GRID, width=1.2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(width); return c


def header(slide, section, heading, n):
    tx(slide, section.upper(), .58, .28, 4.8, .25, 10, BLUE, True)
    tx(slide, heading, .55, .62, 11.7, .58, 27, NAVY, True)
    tx(slide, f"{n:02d}", 12.25, .32, .5, .25, 10, MUTED, True, PP_ALIGN.RIGHT)
    ln(slide, .58, 1.35, 12.75, 1.35, GRID, .8)


def foot(slide, n):
    ln(slide, .58, 7.12, 12.75, 7.12, GRID, .7)
    tx(slide, "Learning to Use References for Skill Induction", .6, 7.19, 5.5, .17, 8, MUTED)
    tx(slide, str(n), 12.2, 7.19, .42, .17, 8, MUTED, False, PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body, fill=WHITE, accent=BLUE, size=16):
    shape(slide, x, y, w, h, fill, GRID, True)
    shape(slide, x, y, .08, h, accent, accent, True)
    tx(slide, title, x+.25, y+.18, w-.45, .32, 16, NAVY, True)
    tx(slide, body, x+.25, y+.62, w-.45, h-.78, size, INK)


def pill(slide, s, x, y, w, fill, color=INK):
    shape(slide, x, y, w, .34, fill, fill, True)
    tx(slide, s, x+.04, y+.04, w-.08, .23, 10, color, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def arrow(slide, x1, y, x2, color=BLUE):
    ln(slide, x1, y, x2, y, color, 2.0)
    tx(slide, ">", (x1+x2)/2-.1, y-.17, .2, .3, 17, color, True, PP_ALIGN.CENTER)


def workflow(slide, labels, x, y, widths, fills):
    cur = x
    for i, label in enumerate(labels):
        shape(slide, cur, y, widths[i], .8, fills[i], fills[i], True)
        tx(slide, label, cur+.06, y+.14, widths[i]-.12, .5, 14, NAVY, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(labels)-1:
            arrow(slide, cur+widths[i]+.03, y+.4, cur+widths[i]+.38)
        cur += widths[i]+.43


def note(slide, s):
    # Speaker notes are provided separately for portability across python-pptx versions.
    return None


# 1 Title
s = prs.slides.add_slide(BLANK); shape(s, 0, 0, 13.333, 7.5, NAVY, NAVY)
shape(s, 0, 5.98, 13.333, 1.52, BLUE, BLUE)
tx(s, "RESEARCH PROPOSAL", .72, .72, 3.5, .28, 12, RGBColor(168, 211, 248), True)
tx(s, "Learning to Use References\nfor Reusable Skill Induction", .7, 1.45, 10.7, 1.32, 34, WHITE, True)
tx(s, "From trajectory-local lessons to state-conditioned procedures", .75, 3.35, 8.8, .36, 20, RGBColor(218, 232, 248))
tx(s, "Author Name  ·  Affiliation", .75, 6.38, 5.0, .28, 15, WHITE, True)
tx(s, "Working title", 10.0, 6.4, 2.4, .25, 12, RGBColor(218, 232, 248), False, PP_ALIGN.RIGHT)

# 2 Motivation
s = prs.slides.add_slide(BLANK); header(s, "Motivation", "A skill should be reusable — not a transcript in disguise", 2)
tx(s, "Trajectory-to-skill methods produce useful rules, but the rules often remain tied to how a training trace happened to look.", .82, 1.7, 11.7, .5, 21, NAVY, True, PP_ALIGN.CENTER)
workflow(s, ["trajectory\ntrace", "textual\nskill", "new\ntrajectory"], .95, 2.65, [2.0, 2.0, 2.0], [PALE_BLUE, PALE_ORANGE, PALE_TEAL])
tx(s, "Where generalization breaks", .9, 4.1, 3.0, .3, 13, RED, True)
bullet_box(s, ["Concrete entities become fixed rules", "Different workflow states are flattened", "Reference examples are copied, not interpreted"], .9, 4.48, 5.45, 1.25, 17)
shape(s, 7.0, 4.12, 5.25, 1.55, PALE_RED, PALE_RED, True)
tx(s, "Research question", 7.35, 4.42, 2.2, .28, 13, RED, True)
tx(s, "Can a skill learn when to consult\nreference — and how to use it?", 7.35, 4.82, 4.45, .58, 23, NAVY, True)
foot(s, 2)

# 3 Existing generated skills
s = prs.slides.add_slide(BLANK); header(s, "Observation", "Two generated skills reveal the same missing capability", 3)
card(s, .72, 1.75, 5.72, 4.38, "AWM-style workflow", "Rich local experience\n\n• concrete user / credential values\n• many action patterns\n• conversational heuristics\n• weak parameterization and prioritization", PALE_ORANGE, ORANGE, 17)
card(s, 6.9, 1.75, 5.72, 4.38, "Trace2Skill-style protocol", "Precise action-slot constraints\n\n• action-first decision\n• exact slot ordering\n• no stale values\n• state and schema remain implicit", PALE_BLUE, BLUE, 17)
tx(s, "Different symptoms, shared gap: the skill does not operationalize evidence use.", 1.0, 6.45, 11.3, .32, 19, NAVY, True, PP_ALIGN.CENTER)
foot(s, 3)

# 4 Mock case macro
s = prs.slides.add_slide(BLANK); header(s, "Bad case 1 · illustrative", "Macro skill: the workflow is remembered, but not parameterized", 4)
tx(s, "TRAINING TRACE", .82, 1.72, 2.0, .25, 11, ORANGE, True)
shape(s, .75, 2.05, 5.35, 2.65, PALE_ORANGE, PALE_ORANGE, True)
tx(s, "Name: Chloe Zhang\nPhone: (482) 837-8571\nAction: verify-identity\n\nGenerated skill:\nverify-identity:chloe zhang,...", 1.05, 2.35, 4.75, 2.0, 18, INK, False, font="Consolas")
tx(s, "TEST TRACE · MOCK", 6.75, 1.72, 2.0, .25, 11, BLUE, True)
shape(s, 6.65, 2.05, 5.85, 2.65, PALE_BLUE, PALE_BLUE, True)
tx(s, "Name: Alice Smith\nPhone: 51909\nSame subflow, new entity values\n\nExpected:\nverify-identity(Alice Smith, 51909, ...)", 6.98, 2.35, 5.15, 2.0, 18, INK, False, font="Consolas")
pill(s, "BAD CASE TO INSTANTIATE", 4.2, 5.1, 4.2, PALE_RED, RED)
tx(s, "Root cause: a trajectory example was stored as a procedure,\nbut the procedure did not expose variables, slots, or grounding rules.", 1.35, 5.72, 10.7, .7, 19, NAVY, True, PP_ALIGN.CENTER)
foot(s, 4)

# 5 Mock case branch
s = prs.slides.add_slide(BLANK); header(s, "Bad case 2 · illustrative", "Branch collapse: valid rules look contradictory when state is hidden", 5)
tx(s, "TRAJECTORY A", .85, 1.72, 1.8, .25, 11, BLUE, True)
workflow(s, ["account\nidentified", "2 creds\navailable", "verify", "verified"], .8, 2.25, [1.6, 1.45, 1.2, 1.3], [PALE_BLUE, PALE_TEAL, PALE_BLUE, PALE_TEAL])
tx(s, "TRAJECTORY B", .85, 3.72, 1.8, .25, 11, ORANGE, True)
workflow(s, ["account\nidentified", "1 cred\navailable", "ask for\ncredential", "verify"], .8, 4.25, [1.6, 1.45, 1.35, 1.2], [PALE_BLUE, PALE_ORANGE, PALE_ORANGE, PALE_BLUE])
shape(s, 8.0, 2.0, 4.45, 3.15, PALE_RED, PALE_RED, True)
tx(s, "FLATTENED SKILL", 8.35, 2.3, 2.3, .25, 11, RED, True)
tx(s, "“Verify identity\nwith credentials.”", 8.35, 2.82, 3.7, .8, 25, NAVY, True, PP_ALIGN.CENTER)
tx(s, "Which rule wins?\nThe missing object is the state condition.", 8.35, 4.05, 3.7, .65, 18, RED, True, PP_ALIGN.CENTER)
pill(s, "BAD CASE TO INSTANTIATE", 4.2, 5.88, 4.2, PALE_RED, RED)
foot(s, 5)

# 6 Insight
s = prs.slides.add_slide(BLANK); header(s, "Insight", "Reference should be an operational resource, not passive context", 6)
tx(s, "Current pattern", 1.0, 1.75, 2.0, .28, 12, RED, True)
workflow(s, ["query by\nkeywords", "paste\nsnippets", "generate\naction"], 1.0, 2.35, [1.85, 1.85, 1.85], [PALE_RED, PALE_RED, PALE_RED])
tx(s, "Proposed pattern", 7.0, 1.75, 2.2, .28, 12, TEAL, True)
workflow(s, ["estimate\nstate", "decide\nretrieve", "query\nneed", "ground +\nverify", "act"], 6.8, 2.35, [1.45, 1.45, 1.35, 1.55, 1.0], [PALE_BLUE, PALE_ORANGE, PALE_TEAL, PALE_BLUE, PALE_TEAL])
shape(s, 2.0, 4.55, 9.3, 1.25, NAVY, NAVY, True)
tx(s, "The skill must learn a reference-use policy:\nwhen to retrieve · what to ask for · how to apply · when to verify", 2.35, 4.84, 8.6, .65, 22, WHITE, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
foot(s, 6)

# 7 representation
s = prs.slides.add_slide(BLANK); header(s, "Representation", "Skill = task policy + reference-use policy", 7)
card(s, .75, 1.85, 5.55, 3.95, "Task policy", "State → business action\n\naccount_identified + enough credentials\n→ verify-identity\n\nidentity_verified\n→ send-link", PALE_BLUE, BLUE, 18)
card(s, 7.0, 1.85, 5.55, 3.95, "Reference-use policy", "State + uncertainty → query → evidence use\n\nRetrieve schema, branch, or boundary case.\nGround values in current dialogue.\nValidate applicability before acting.", PALE_TEAL, TEAL, 18)
tx(s, "A reusable skill is not a longer prompt. It is a policy for acting and consulting evidence.", 1.0, 6.3, 11.3, .36, 19, NAVY, True, PP_ALIGN.CENTER)
foot(s, 7)

# 8 method overview
s = prs.slides.add_slide(BLANK); header(s, "Method", "Learn to retrieve, ground, apply, and verify", 8)
workflow(s, ["trajectory\nstate", "retrieve\nor act", "state-aware\nquery", "reference\nevidence", "ground\ncurrent slots", "execute"], .62, 2.25, [1.45, 1.35, 1.5, 1.55, 1.55, 1.2], [PALE_BLUE, PALE_ORANGE, PALE_TEAL, PALE_BLUE, PALE_ORANGE, PALE_TEAL])
for i, (head, body, x, fill, accent) in enumerate([
    ("Access", "Should I consult\nreference first?", .85, PALE_ORANGE, ORANGE),
    ("Query", "What uncertainty\nneeds resolving?", 4.25, PALE_TEAL, TEAL),
    ("Use", "Which constraints\napply to this state?", 7.65, PALE_BLUE, BLUE),
    ("Verify", "Are action and slots\ngrounded now?", 10.15, PALE_RED, RED),
]):
    shape(s, x, 4.25, 2.25 if i < 3 else 2.2, 1.25, fill, fill, True)
    tx(s, head, x+.15, 4.48, 1.95, .25, 14, accent, True, PP_ALIGN.CENTER)
    tx(s, body, x+.12, 4.82, 2.0, .5, 16, NAVY, True, PP_ALIGN.CENTER)
tx(s, "Reference access becomes part of the induced procedure.", 1.0, 6.25, 11.3, .35, 20, NAVY, True, PP_ALIGN.CENTER)
foot(s, 8)

# 9 learning
s = prs.slides.add_slide(BLANK); header(s, "Learning", "The training signal is a decision trace, not only an action trace", 9)
shape(s, .78, 1.82, 11.75, 1.18, PALE_BLUE, PALE_BLUE, True)
tx(s, "state → query decision → retrieved evidence → selected rule → action → outcome", 1.05, 2.14, 11.2, .42, 22, NAVY, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
card(s, .78, 3.5, 3.55, 1.85, "Supervise access", "Retrieve vs. act\nwhen uncertainty is high", PALE_ORANGE, ORANGE, 16)
card(s, 4.88, 3.5, 3.55, 1.85, "Supervise query", "Schema vs. branch\nvs. boundary evidence", PALE_TEAL, TEAL, 16)
card(s, 8.98, 3.5, 3.55, 1.85, "Supervise use", "Evidence must change\n(or validate) the action", PALE_BLUE, BLUE, 16)
tx(s, "Learning objective: improve decisions per retrieved token — not retrieve more text.", 1.0, 6.25, 11.3, .35, 19, NAVY, True, PP_ALIGN.CENTER)
foot(s, 9)

# 10 experiments
s = prs.slides.add_slide(BLANK); header(s, "Evaluation", "The paper needs three controlled bad-case families", 10)
rows = [
    ("Macro", "seen workflow / unseen entities", "parameterized generalization", PALE_ORANGE, ORANGE),
    ("Branch", "same intent / different state", "conditional action accuracy", PALE_TEAL, TEAL),
    ("Reference", "retrieved example conflicts with dialogue", "grounding and misuse rate", PALE_BLUE, BLUE),
]
for i, (a,b,c,fill,accent) in enumerate(rows):
    y = 1.85 + i*1.22
    shape(s, .8, y, 2.0, .85, fill, fill, True); tx(s, a, .95, y+.23, 1.7, .3, 17, accent, True, PP_ALIGN.CENTER)
    shape(s, 3.05, y, 3.85, .85, LIGHT, GRID, True); tx(s, b, 3.25, y+.23, 3.45, .3, 16, NAVY, True, PP_ALIGN.CENTER)
    shape(s, 7.2, y, 5.25, .85, WHITE, GRID, True); tx(s, c, 7.45, y+.23, 4.75, .3, 16, INK, True, PP_ALIGN.CENTER)
pill(s, "MOCK / TO BE REPLACED BY ALIGNED CASES", 3.5, 5.78, 6.3, PALE_RED, RED)
tx(s, "A valid claim requires: original trace + generated skill + retrieval context + prediction + outcome.", 1.0, 6.35, 11.3, .3, 15, MUTED, False, PP_ALIGN.CENTER)
foot(s, 10)

# 11 contributions
s = prs.slides.add_slide(BLANK); header(s, "Conclusion", "The contribution is a skill that knows how to consult evidence", 11)
shape(s, .8, 1.8, 11.75, 1.35, NAVY, NAVY, True)
tx(s, "Existing methods learn what to do.\nWe learn when and how to consult evidence before doing it.", 1.1, 2.08, 11.15, .75, 25, WHITE, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
bullet_box(s, ["Macro skill: parameterized, reusable procedures", "Branch-aware skill: conditions separate valid rules", "Reference-use policy: retrieve, ground, apply, verify", "Evaluation: bad case → cause → insight → measurable fix"], 1.05, 3.75, 8.7, 1.8, 19)
shape(s, 9.6, 4.0, 2.35, 1.5, PALE_TEAL, PALE_TEAL, True)
tx(s, "Reusable\n+\nState-aware", 9.8, 4.28, 1.95, .9, 22, TEAL, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
tx(s, "Thank you", 10.1, 6.28, 2.3, .35, 20, BLUE, True, PP_ALIGN.CENTER)
foot(s, 11)


prs.save(OUT)
print(f"saved {OUT}")
